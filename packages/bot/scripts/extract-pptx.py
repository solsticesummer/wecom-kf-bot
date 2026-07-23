#!/usr/bin/env python3
"""Extract text from a .pptx, one record per slide, as JSON to stdout.

Usage:  python3 scripts/extract-pptx.py <file.pptx>
Output: JSON array of {"index": int, "title": str, "text": str}

build-index.js shells out to this and embeds one chunk per slide whose `text` is
non-empty (image-only slides come back with empty text and are skipped). Text only —
we do not extract slide images in Phase 1.
"""
import sys
import json
from pptx import Presentation


def slide_records(path):
    prs = Presentation(path)
    out = []
    for i, slide in enumerate(prs.slides):
        title_shape = slide.shapes.title
        title = title_shape.text.strip() if title_shape is not None else ""

        parts = []
        for shape in slide.shapes:
            if shape is title_shape:
                continue  # already captured as the title
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    parts.append(t)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    parts.append(" | ".join(cells))

        out.append({"index": i, "title": title, "text": "\n".join(parts).strip()})
    return out


if __name__ == "__main__":
    if len(sys.argv) != 2:
        sys.exit("usage: extract-pptx.py <file.pptx>")
    json.dump(slide_records(sys.argv[1]), sys.stdout, ensure_ascii=False)
