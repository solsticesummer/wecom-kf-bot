"""Source connectors: turn a file or project folder into chunked documents.

Handles local text (markdown/text via heading-aware chunking, code via line-window
chunking) plus the two binary formats the bot's own corpus uses: PDF and PPTX. Everything
plugs in behind the same ``collect(path) -> list[doc]`` shape, so a folder holding a mix of
docs, slides and code ingests in one pass.

The binary extractors mirror the TS bot's build-index sources so the two stores see the same
text: ``pdftotext -layout`` (poppler) for PDF, python-pptx one-record-per-slide for PPTX.
Both are OPTIONAL dependencies — a missing tool skips those files with a warning rather than
failing the whole ingest, since most projects have neither.
"""

from __future__ import annotations

import os
import shutil
import subprocess

from . import chunk

TEXT_EXT = {".md", ".markdown", ".txt", ".rst"}
CODE_EXT = {
    ".py", ".ts", ".tsx", ".js", ".jsx", ".mjs", ".cjs", ".json", ".yaml", ".yml",
    ".toml", ".sql", ".sh", ".html", ".css", ".go", ".rs", ".java", ".rb", ".php",
}
SKIP_DIRS = {".git", "node_modules", ".venv", "venv", "dist", "build", "__pycache__",
             ".next", ".turbo", "coverage", ".mypy_cache", ".pytest_cache"}


PDF_EXT = {".pdf"}
PPTX_EXT = {".pptx"}


def _pdf_pieces(path: str) -> list[tuple[str, str]]:
    """PDF → chunks via ``pdftotext -layout`` (poppler).

    ``-layout`` preserves column/table geometry, which is the whole point: without it a
    pricing table collapses into interleaved prose and the numbers stop lining up with
    their labels. Extracted text has no markdown headings, so it goes through the
    line-window splitter.
    """
    if not shutil.which("pdftotext"):
        print(f"skip {os.path.basename(path)}: pdftotext not installed (brew install poppler)")
        return []
    try:
        text = subprocess.run(
            ["pdftotext", "-layout", path, "-"],
            capture_output=True,
            check=True,
            timeout=120,
        ).stdout.decode("utf-8", errors="replace")
    except (subprocess.CalledProcessError, subprocess.TimeoutExpired) as exc:
        print(f"skip {os.path.basename(path)}: pdftotext failed ({exc})")
        return []
    name = os.path.basename(path)
    return [(f"{name}#{i}", c) for i, c in enumerate(chunk.chunk_text(text))]


def _pptx_pieces(path: str) -> list[tuple[str, str]]:
    """PPTX → one chunk per slide, titled by the slide's title placeholder.

    A slide is already a human-authored unit of meaning, so it is a better chunk boundary
    than any window we could impose. Image-only slides yield no text and are skipped.
    """
    try:
        from pptx import Presentation  # optional dependency
    except ImportError:
        print(f"skip {os.path.basename(path)}: python-pptx not installed (pip install python-pptx)")
        return []

    pieces: list[tuple[str, str]] = []
    prs = Presentation(path)
    for i, slide in enumerate(prs.slides):
        title_shape = slide.shapes.title
        title = title_shape.text.strip() if title_shape is not None else ""
        parts: list[str] = []
        for shape in slide.shapes:
            if shape is title_shape:
                continue  # already captured as the title
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text.strip())
            if shape.has_table:
                for row in shape.table.rows:
                    parts.append(" | ".join(c.text.strip() for c in row.cells))
        body = "\n".join(parts).strip()
        if not (title or body):
            continue
        pieces.append((title or f"slide {i + 1}", (f"{title}\n" if title else "") + body))
    return pieces


def collect(root: str) -> list[dict]:
    """Walk ``root`` and return ``[{source, section, content}, ...]``.

    ``source`` is the path relative to ``root`` so results are portable.
    """
    root = os.path.abspath(root)
    if os.path.isfile(root):
        files = [root]
        base = os.path.dirname(root)
    else:
        files = []
        base = root
        for dirpath, dirnames, filenames in os.walk(root):
            dirnames[:] = [d for d in dirnames if d not in SKIP_DIRS]
            for fn in filenames:
                files.append(os.path.join(dirpath, fn))

    docs: list[dict] = []
    for path in files:
        ext = os.path.splitext(path)[1].lower()
        rel = os.path.relpath(path, base)
        pieces: list[tuple[str, str]] = []

        # Binary formats first: they need an external extractor, not a text read.
        if ext in PDF_EXT:
            pieces = _pdf_pieces(path)
        elif ext in PPTX_EXT:
            pieces = _pptx_pieces(path)
        elif ext in TEXT_EXT or ext in CODE_EXT:
            try:
                with open(path, encoding="utf-8") as fh:
                    text = fh.read()
            except (UnicodeDecodeError, OSError):
                continue  # binary or unreadable — skip
            if not text.strip():
                continue
            if ext in TEXT_EXT:
                pieces = chunk.chunk_markdown(text)
            if not pieces:  # non-markdown, or markdown with no ## headings
                pieces = [(f"{rel}#{i}", c) for i, c in enumerate(chunk.chunk_text(text))]
        else:
            continue

        for section, content in pieces:
            docs.append({"source": rel, "section": section, "content": content})
    return docs
